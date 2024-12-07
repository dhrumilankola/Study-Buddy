import os
from pptx import Presentation
from nbformat import read
from langchain_community.document_loaders import PyPDFLoader
from langchain.schema import Document

def extract_text_from_pdf(pdf_path):
    pdf_loader = PyPDFLoader(pdf_path)
    return [Document(page_content=doc.page_content, metadata={"source": pdf_path}) for doc in pdf_loader.load()]

def extract_text_from_pptx(pptx_path):
    presentation = Presentation(pptx_path)
    pptx_text = ""
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                pptx_text += shape.text + "\n\n"
    return [Document(page_content=pptx_text, metadata={"source": pptx_path})]

def extract_text_from_notebook(notebook_path):
    with open(notebook_path, 'r', encoding='utf-8') as f:
        notebook = read(f, as_version=4)
    notebook_text = ""
    for cell in notebook.cells:
        if cell.cell_type == 'markdown' or cell.cell_type == 'code':
            notebook_text += cell.source + "\n\n"
    return [Document(page_content=notebook_text, metadata={"source": notebook_path})]

def process_file(file_path):
    extension = os.path.splitext(file_path)[1].lower()
    if extension == ".pdf":
        return extract_text_from_pdf(file_path)
    elif extension == ".pptx":
        return extract_text_from_pptx(file_path)
    elif extension == ".ipynb":
        return extract_text_from_notebook(file_path)
    else:
        raise ValueError(f"Unsupported file type: {extension}")
