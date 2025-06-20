from typing import List, Dict, Optional, Tuple, Any
import PyPDF2
from pptx import Presentation
import nbformat
import os
import json
from datetime import datetime
import re
from langchain_text_splitters import RecursiveCharacterTextSplitter, MarkdownTextSplitter
from app.models.schemas import Document
import logging
import asyncio
from langchain_community.document_loaders import (
    PyPDFLoader, 
    TextLoader,
    UnstructuredPowerPointLoader,
    NotebookLoader
)
import tempfile

logger = logging.getLogger(__name__)

class EnhancedDocumentProcessor:
    """Advanced document processor with better text extraction and chunking"""
    
    def __init__(self):
        self.supported_extensions = set(os.environ.get('ALLOWED_EXTENSIONS', 
                                                      '.pdf,.txt,.pptx,.ipynb').split(','))
        self.chunk_size = int(os.environ.get('CHUNK_SIZE', '1000'))
        self.chunk_overlap = int(os.environ.get('CHUNK_OVERLAP', '200'))

    async def process_file(self, file_path: str) -> List[str]:
        """Process file using LangChain document loaders for better extraction"""
        file_extension = os.path.splitext(file_path)[1].lower()
        
        if file_extension not in self.supported_extensions:
            raise ValueError(f"Unsupported file type: {file_extension}")
            
        try:
            # Use LangChain loaders for better extraction
            if file_extension == '.pdf':
                return await self._process_with_langchain_loader(PyPDFLoader, file_path)
            elif file_extension == '.txt':
                return await self._process_with_langchain_loader(TextLoader, file_path)
            elif file_extension == '.pptx':
                return await self._process_with_langchain_loader(UnstructuredPowerPointLoader, file_path)
            elif file_extension == '.ipynb':
                return await self._process_with_langchain_loader(NotebookLoader, file_path)
            else:
                raise ValueError(f"Unsupported file type: {file_extension}")
        except Exception as e:
            logger.error(f"Error processing file {file_path}: {str(e)}")
            # Fall back to traditional processing methods
            logger.info(f"Falling back to traditional processing for {file_path}")
            return await self._fallback_processing(file_path)
            
    async def _process_with_langchain_loader(self, loader_class, file_path: str) -> List[str]:
        """Process a file using a LangChain document loader"""
        try:
            loader = loader_class(file_path)
            # Run in a thread pool to avoid blocking
            docs = await asyncio.to_thread(loader.load)
            
            # Extract text from documents
            texts = [doc.page_content for doc in docs]
            return texts
        except Exception as e:
            logger.error(f"Error with LangChain loader: {str(e)}")
            raise
            
    async def _fallback_processing(self, file_path: str) -> List[str]:
        """Fallback to traditional processing methods if LangChain loaders fail"""
        file_extension = os.path.splitext(file_path)[1].lower()
        
        if file_extension == '.pdf':
            return await asyncio.to_thread(self._process_pdf_traditional, file_path)
        elif file_extension == '.txt':
            return await asyncio.to_thread(self._process_txt_traditional, file_path)
        elif file_extension == '.pptx':
            return await asyncio.to_thread(self._process_pptx_traditional, file_path)
        elif file_extension == '.ipynb':
            return await asyncio.to_thread(self._process_ipynb_traditional, file_path)
        else:
            raise ValueError(f"Unsupported file type: {file_extension}")
            
    def _process_pdf_traditional(self, file_path: str) -> List[str]:
        """Extract text from PDF file using PyPDF2"""
        texts = []
        try:
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                for page in pdf_reader.pages:
                    text = page.extract_text()
                    if text.strip():
                        texts.append(text)
        except Exception as e:
            logger.error(f"Error processing PDF {file_path}: {str(e)}")
            raise
        return texts

    def _process_pptx_traditional(self, file_path: str) -> List[str]:
        """Extract text from PPTX file"""
        texts = []
        try:
            presentation = Presentation(file_path)
            for slide in presentation.slides:
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slide_text.append(shape.text)
                if slide_text:
                    texts.append("\n".join(slide_text))
        except Exception as e:
            logger.error(f"Error processing PPTX {file_path}: {str(e)}")
            raise
        return texts

    def _process_txt_traditional(self, file_path: str) -> List[str]:
        """Process text file"""
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                text = file.read()
                # Return as a single document for chunking later
                return [text]
        except UnicodeDecodeError:
            # Try with different encodings if UTF-8 fails
            try:
                with open(file_path, 'r', encoding='latin-1') as file:
                    text = file.read()
                    return [text]
            except Exception as e:
                logger.error(f"Error processing TXT with latin-1 {file_path}: {str(e)}")
                raise
        except Exception as e:
            logger.error(f"Error processing TXT {file_path}: {str(e)}")
            raise

    def _process_ipynb_traditional(self, file_path: str) -> List[str]:
        """Process Jupyter notebook"""
        texts = []
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                notebook = nbformat.read(file, as_version=4)
                for cell in notebook.cells:
                    if cell.cell_type == 'markdown':
                        texts.append(cell.source)
                    elif cell.cell_type == 'code':
                        # Include both code and outputs
                        code_text = f"```python\n{cell.source}\n```"
                        texts.append(code_text)
                        if 'outputs' in cell and cell.outputs:
                            output_texts = []
                            for output in cell.outputs:
                                if 'text' in output:
                                    output_texts.append(output.text)
                                elif 'data' in output and 'text/plain' in output.data:
                                    output_texts.append(output.data['text/plain'])
                            if output_texts:
                                texts.append("Output:\n" + "\n".join(output_texts))
        except Exception as e:
            logger.error(f"Error processing Jupyter notebook {file_path}: {str(e)}")
            raise
        return texts

    def chunk_texts(self, texts: List[str], file_type: str) -> List[str]:
        """Chunk texts based on file type with appropriate splitters"""
        if not texts:
            return []
            
        # Join all texts together for consistent chunking
        all_text = "\n\n".join(texts)
        
        # Select appropriate splitter based on file type
        if file_type == 'ipynb':
            # Use markdown splitter for notebooks
            splitter = MarkdownTextSplitter(
                chunk_size=self.chunk_size,
                chunk_overlap=self.chunk_overlap
            )
        else:
            # Use recursive character splitter for other types
            splitter = RecursiveCharacterTextSplitter(
                chunk_size=self.chunk_size,
                chunk_overlap=self.chunk_overlap,
                separators=["\n\n", "\n", ". ", " ", ""]
            )
            
        # Split the combined text
        chunks = splitter.split_text(all_text)
        
        # Clean and filter chunks
        processed_chunks = []
        for chunk in chunks:
            # Clean up whitespace
            cleaned = re.sub(r'\s+', ' ', chunk).strip()
            if len(cleaned) > 10:  # Only keep chunks with meaningful content
                processed_chunks.append(cleaned)
                
        return processed_chunks

    async def process_document(self, document: Document, file_path: str) -> List[Dict]:
        """Process document and prepare for vector store"""
        try:
            file_extension = os.path.splitext(file_path)[1].lower()
            
            if file_extension not in self.supported_extensions:
                raise ValueError(f"Unsupported file type: {file_extension}")

            # Process file to extract text
            extracted_texts = await self.process_file(file_path)
            
            # Chunk the extracted texts
            chunked_texts = self.chunk_texts(extracted_texts, document.file_type)
            
            # Prepare metadata for each chunk
            processed_chunks = []
            for idx, text in enumerate(chunked_texts):
                metadata = {
                    "document_id": document.id,  # Keep as string to match schema
                    "uuid_filename": document.uuid_filename if hasattr(document, 'uuid_filename') and document.uuid_filename else f"{document.id}.{document.file_type}",
                    "filename": document.filename,
                    "chunk_index": idx,
                    "file_type": document.file_type,
                    "processed_date": datetime.now().isoformat(),
                    "total_chunks": len(chunked_texts)
                }
                processed_chunks.append({
                    "text": text,
                    "metadata": metadata
                })
            
            logger.info(f"Document {document.filename} processed into {len(processed_chunks)} chunks")
            return processed_chunks

        except Exception as e:
            logger.error(f"Error processing document {file_path}: {str(e)}")
            raise